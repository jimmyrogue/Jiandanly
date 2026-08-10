"""PowerPoint-specific office.* read and write tools."""

from __future__ import annotations

import io
import os
from pathlib import Path
from typing import Any

from langchain_core.tools import tool
from pptx import Presentation as PptxPresentation
from pptx.util import Inches

from .office_common import (
    atomic_write as _atomic_write,
)
from .office_common import (
    ensure_copy_for_write as _ensure_copy_for_write,
)
from .office_common import (
    resolve_read_source as _resolve_read_source,
)
from .office_common import (
    validate_path as _validate_path,
)
from .office_common import (
    verify_file as _verify_file,
)
from .office_common import (
    write_error as _write_error,
)
from .office_common import (
    write_result as _write_result,
)


def _slide_title_text(slide) -> str:
    """Extract the title text from a slide, or "" if no title placeholder."""
    # python-pptx exposes a `.shapes.title` shortcut, but it's None
    # when the slide layout has no title placeholder. Fall back to the
    # first text-bearing shape so we still surface SOMETHING for
    # title-less slides (cover layouts use centered text frames).
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        text = slide.shapes.title.text_frame.text.strip()
        if text:
            return text
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text.splitlines()[0]
    return ""


def _slide_bullets(slide) -> list[str]:
    """Extract paragraph-level text from the slide's content placeholder.

    A slide typically has one title placeholder + one body placeholder.
    We pull text from every non-title shape's text_frame and split by
    paragraph. Empty paragraphs are dropped.
    """
    bullets: list[str] = []
    title_shape = slide.shapes.title
    # python-pptx returns different wrapper instances for the same
    # underlying XML element across `shapes.title` vs iteration, so an
    # `is` check fails. Compare the underlying lxml elements instead.
    title_elem = title_shape._element if title_shape is not None else None
    for shape in slide.shapes:
        if title_elem is not None and shape._element is title_elem:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                bullets.append(text)
    return bullets


def _slide_notes(slide) -> str:
    """Speaker notes for a slide; "" when no notes slide is attached."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _outline_pptx(path: Any) -> dict[str, Any]:
    """Structural summary of a .pptx — per-slide title, bullets, notes,
    shape counts. Both `office.outline` and `office.read_slides` consume
    this; the panel preview also fetches it through a small HTTP
    endpoint."""
    prs = PptxPresentation(path)
    slides: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        slides.append(
            {
                "index": idx,
                "layout": slide.slide_layout.name if slide.slide_layout else "",
                "title": _slide_title_text(slide),
                "bullets": _slide_bullets(slide),
                "notes": _slide_notes(slide),
                "shape_count": len(slide.shapes),
                "image_count": sum(1 for sh in slide.shapes if sh.shape_type == 13),  # 13 = PICTURE
            }
        )
    return {"slides": slides, "slide_count": len(slides)}


# ═══════════════════════════════════════════════════════════════════════
# PPTX helpers
# ═══════════════════════════════════════════════════════════════════════


def _resolve_slide_layout(prs, layout_name: str | None):
    """Return a slide_layout by name; fall back to layout index 1
    ('Title and Content') when name is None or not found."""
    if layout_name:
        for layout in prs.slide_layouts:
            if layout.name == layout_name:
                return layout
    # python-pptx default deck has these layouts in this order:
    #   0 Title Slide   1 Title and Content   2 Section Header
    #   3 Two Content   4 Comparison          5 Title Only
    #   6 Blank         7 Content with Caption 8 Picture with Caption
    return prs.slide_layouts[1]


def _set_placeholder_text(placeholder, value: str) -> None:
    """Replace a placeholder's text frame with `value`, preserving the
    first paragraph's formatting. python-pptx's `placeholder.text = ...`
    setter does this but drops bullet formatting; we want a single
    paragraph rewrite that survives the round-trip."""
    placeholder.text_frame.text = value


def _set_placeholder_bullets(placeholder, bullets: list[str]) -> None:
    """Replace a placeholder's text frame with one paragraph per bullet."""
    tf = placeholder.text_frame
    tf.clear()  # leaves one empty paragraph
    if not bullets:
        return
    # The cleared text_frame has exactly one paragraph; reuse it for
    # the first bullet, then add the rest.
    tf.paragraphs[0].text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet


def _find_content_placeholder(slide):
    """Return the body placeholder (anything not the title). None when
    the layout has no body slot — caller should handle by adding a
    text box or returning an error."""
    title_shape = slide.shapes.title
    title_elem = title_shape._element if title_shape is not None else None
    for shape in slide.placeholders:
        if title_elem is not None and shape._element is title_elem:
            continue
        return shape
    return None


def _set_slide_content(slide, title: str | None, bullets: list[str] | None) -> None:
    """Apply title/bullets to a slide, leaving unset fields alone."""
    if title is not None:
        if slide.shapes.title is None:
            raise ValueError("slide layout has no title placeholder")
        _set_placeholder_text(slide.shapes.title, title)
    if bullets is not None:
        body = _find_content_placeholder(slide)
        if body is None:
            raise ValueError("slide layout has no body placeholder for bullets")
        _set_placeholder_bullets(body, bullets)


def _validate_slide_index(prs, index: int) -> None:
    n = len(prs.slides)
    if not (0 <= index < n):
        raise ValueError(f"slide index {index} out of range (deck has {n} slides)")


# ═══════════════════════════════════════════════════════════════════════
# PPTX write tools
# ═══════════════════════════════════════════════════════════════════════


@tool("office.create_pptx")
def office_create_pptx(path: str, title: str | None = None) -> dict[str, Any]:
    """Create a new blank PowerPoint deck at `path` with one title slide.

    Special-case among write tools: this writes DIRECTLY to `path`
    (not to a `.edited.pptx` copy) — there's no pre-existing original
    to protect when you're creating from scratch. Subsequent edits
    via office.add_slide / office.update_slide / etc. follow the
    standard copy-on-first-write pattern and land in
    `<basename>.edited.pptx`.

    Args:
        path: absolute path where the new .pptx should land (extension
              must be .pptx).
        title: optional title for the first slide. When omitted the
               slide is created with an empty title placeholder.

    Returns:
        dict with `ok`, `original_path`, `edited_path` (= path,
        because we wrote the original), `kind="powerpoint"`,
        `summary={slide_count: 1}`.
    """
    if not path:
        return _write_error(None, "path required")
    resolved = os.path.abspath(os.path.expanduser(path))
    if Path(resolved).suffix.lower() != ".pptx":
        return _write_error(None, "office.create_pptx requires a .pptx path")
    if os.path.exists(resolved):
        return _write_error(
            "powerpoint",
            f"file already exists: {resolved} (use office.update_slide to edit it)",
            resolved,
        )
    try:
        prs = PptxPresentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        if title is not None and slide.shapes.title is not None:
            _set_placeholder_text(slide.shapes.title, title)
        # Ensure the parent directory exists (project-relative paths).
        Path(resolved).parent.mkdir(parents=True, exist_ok=True)
        prs.save(resolved)
        _verify_file(resolved, "powerpoint")
    except Exception as exc:
        # Clean up partial write so a half-created file doesn't linger.
        if os.path.exists(resolved):
            try:
                os.remove(resolved)
            except OSError:
                pass
        return _write_error(
            "powerpoint", f"create failed: {exc.__class__.__name__}: {exc}", resolved
        )
    return {
        "ok": "true",
        "original_path": resolved,
        "edited_path": resolved,  # the create IS the working file
        "kind": "powerpoint",
        "summary": {"slide_count": 1, "title": title or ""},
    }


@tool("office.add_slide")
def office_add_slide(
    path: str,
    layout: str | None = None,
    title: str | None = None,
    bullets: list[str] | None = None,
) -> dict[str, Any]:
    """Append a new slide to a .pptx. Writes to `<basename>.edited.pptx`.

    Args:
        path: .pptx to edit.
        layout: python-pptx layout name (e.g. "Title and Content",
                "Title Slide", "Section Header", "Two Content", "Blank").
                Default: "Title and Content".
        title: optional title text.
        bullets: optional list of strings; each becomes one paragraph
                 (bullet) in the body placeholder.

    Returns:
        dict with `ok`, `original_path`, `edited_path`, `kind`,
        `summary={index, layout}` where `index` is the new slide's
        0-based position.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.add_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        slide_layout = _resolve_slide_layout(prs, layout)
        slide = prs.slides.add_slide(slide_layout)
        _set_slide_content(slide, title, bullets)
        prs.save(tmp_path)
        return {"index": len(prs.slides) - 1, "layout": slide_layout.name}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.update_slide")
def office_update_slide(
    path: str,
    index: int,
    title: str | None = None,
    bullets: list[str] | None = None,
    layout: str | None = None,
) -> dict[str, Any]:
    """Update an existing slide's title / bullets / layout.

    Writes to `<basename>.edited.pptx`. Any argument left as None is
    untouched. NOTE: changing `layout` re-creates the slide on the new
    layout, so any custom shapes you added previously are lost — only
    the title + body bullets carry over.

    Args:
        path: .pptx to edit.
        index: 0-based slide index.
        title, bullets, layout: see add_slide.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.update_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        if layout is not None:
            # Layout change: remove the existing slide, insert a new one
            # at the same position with the new layout.
            existing = prs.slides[index]
            existing_title = _slide_title_text(existing)
            existing_bullets = _slide_bullets(existing)
            new_title = title if title is not None else existing_title
            new_bullets = bullets if bullets is not None else existing_bullets
            new_layout = _resolve_slide_layout(prs, layout)
            new_slide = prs.slides.add_slide(new_layout)
            _set_slide_content(new_slide, new_title, new_bullets)
            # Move the new slide (currently last) to the original index,
            # then delete the original.
            _move_slide(prs, len(prs.slides) - 1, index + 1)  # insert after the original
            _delete_slide(prs, index)
            applied_layout = new_layout.name
        else:
            slide = prs.slides[index]
            _set_slide_content(slide, title, bullets)
            applied_layout = slide.slide_layout.name if slide.slide_layout else ""
        prs.save(tmp_path)
        return {"index": index, "layout": applied_layout}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


def _delete_slide(prs, index: int) -> None:
    """python-pptx has no public delete; manipulate sldIdLst directly."""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    to_remove = slides[index]
    # Drop the relationship so PowerPoint doesn't complain about a
    # dangling rId on open.
    rels_attr = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    rid = to_remove.attrib[rels_attr]
    prs.part.drop_rel(rid)
    sld_id_lst.remove(to_remove)


def _move_slide(prs, from_index: int, to_index: int) -> None:
    """Reorder slides via sldIdLst manipulation. to_index is the
    desired final position; negative indices not supported."""
    sld_id_lst = prs.slides._sldIdLst
    slides = list(sld_id_lst)
    if not (0 <= from_index < len(slides)):
        raise ValueError(f"from_index {from_index} out of range")
    if not (0 <= to_index < len(slides)):
        raise ValueError(f"to_index {to_index} out of range")
    moving = slides[from_index]
    sld_id_lst.remove(moving)
    sld_id_lst.insert(to_index, moving)


@tool("office.delete_slide")
def office_delete_slide(path: str, index: int) -> dict[str, Any]:
    """Delete the slide at `index` (0-based). Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.delete_slide requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        deleted_title = _slide_title_text(prs.slides[index])
        _delete_slide(prs, index)
        prs.save(tmp_path)
        return {
            "deleted_index": index,
            "deleted_title": deleted_title,
            "slide_count": len(prs.slides),
        }

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.reorder_slides")
def office_reorder_slides(path: str, from_index: int, to_index: int) -> dict[str, Any]:
    """Move slide from `from_index` to `to_index` (both 0-based).
    Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.reorder_slides requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _move_slide(prs, from_index, to_index)
        prs.save(tmp_path)
        return {"from_index": from_index, "to_index": to_index}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.set_slide_title")
def office_set_slide_title(path: str, index: int, title: str) -> dict[str, Any]:
    """Set just the title text of slide `index`. Writes to `.edited.pptx`."""
    return office_update_slide.func(path=path, index=index, title=title)


@tool("office.set_slide_bullets")
def office_set_slide_bullets(path: str, index: int, bullets: list[str]) -> dict[str, Any]:
    """Replace the body bullets of slide `index`. Writes to `.edited.pptx`.

    `bullets` is a list of strings — each becomes one paragraph in the
    body placeholder. Pass `[]` to clear the body."""
    return office_update_slide.func(path=path, index=index, bullets=bullets)


@tool("office.set_slide_notes")
def office_set_slide_notes(path: str, index: int, notes: str) -> dict[str, Any]:
    """Set the presenter notes for slide `index`. Writes to `.edited.pptx`."""
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.set_slide_notes requires a .pptx file", resolved)
    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        slide = prs.slides[index]
        # has_notes_slide is True if notes have ever been set; accessing
        # notes_slide on a slide-without-notes auto-creates the notes slide.
        slide.notes_slide.notes_text_frame.text = notes
        prs.save(tmp_path)
        return {"index": index, "notes_excerpt": notes[:60]}

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


@tool("office.add_image_to_slide")
def office_add_image_to_slide(
    path: str,
    index: int,
    image_path: str,
    left_in: float | None = None,
    top_in: float | None = None,
    width_in: float | None = None,
) -> dict[str, Any]:
    """Insert an image into slide `index`. Writes to `.edited.pptx`.

    Args:
        path: .pptx to edit.
        index: 0-based slide index.
        image_path: absolute path to an image file (.png, .jpg, .gif).
                    Must already exist on disk.
        left_in, top_in, width_in: position + size in inches. When
                                   omitted the image is centered on
                                   the slide at ~50% width.
    """
    resolved, kind, err = _validate_path(path)
    if err is not None:
        return _write_error(None, err)
    assert resolved is not None and kind is not None
    if kind != "powerpoint":
        return _write_error(kind, "office.add_image_to_slide requires a .pptx file", resolved)
    img_resolved = os.path.abspath(os.path.expanduser(image_path))
    if not os.path.isfile(img_resolved):
        return _write_error(kind, f"image file not found: {img_resolved}", resolved)

    target = _ensure_copy_for_write(resolved)

    def _do_write(tmp_path: str) -> dict[str, Any]:
        prs = PptxPresentation(target)
        _validate_slide_index(prs, index)
        slide = prs.slides[index]
        # Slide width is usually 10in; height ~7.5in. Default: 5in wide,
        # centered horizontally, ~30% from top.
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        w = Inches(width_in) if width_in is not None else Inches(5)
        left = Inches(left_in) if left_in is not None else (slide_w - w) // 2
        top = Inches(top_in) if top_in is not None else slide_h // 4
        pic = slide.shapes.add_picture(img_resolved, left, top, width=w)
        prs.save(tmp_path)
        return {
            "index": index,
            "image": img_resolved,
            "width_in": float(pic.width) / 914400,  # EMU → inches
        }

    try:
        summary = _atomic_write(target, kind, _do_write)
    except Exception as exc:
        return _write_error(kind, f"write failed: {exc.__class__.__name__}: {exc}", resolved)
    return _write_result(resolved, target, kind, summary)


# ═══════════════════════════════════════════════════════════════════════
# PPTX read tool
# ═══════════════════════════════════════════════════════════════════════


@tool("office.read_slides")
def office_read_slides(path: str) -> dict[str, Any]:
    """Return per-slide structured data (title, bullets, notes, …).

    Use this instead of `office.read` when you want LLM-friendly slide
    semantics (one entry per slide with title + bullets list) rather
    than the compact markdown dump. The frontend preview panel
    also calls this through a small HTTP endpoint to render the
    outline view.

    Args:
        path: Runtime `/attachments/...` route or workspace .pptx path.

    Returns:
        dict with:
          ok, path, kind="powerpoint"
          slides (list of {index, layout, title, bullets, notes,
                  shape_count, image_count})
          slide_count
    """
    source, kind, err = _resolve_read_source(path)
    if err is not None:
        return {"ok": "false", "error": err}
    assert source is not None and kind is not None
    reported_path = path if isinstance(source, bytes) else source
    if kind != "powerpoint":
        return {
            "ok": "false",
            "path": reported_path,
            "error": "office.read_slides requires a .pptx file",
        }
    try:
        details = _outline_pptx(io.BytesIO(source) if isinstance(source, bytes) else source)
    except Exception as exc:
        return {
            "ok": "false",
            "path": reported_path,
            "kind": kind,
            "error": f"failed to read .pptx: {exc.__class__.__name__}: {exc}",
        }
    return {"ok": "true", "path": reported_path, "kind": kind, **details}
