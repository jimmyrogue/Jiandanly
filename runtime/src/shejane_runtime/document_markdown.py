"""Convert model-readable document formats without content-detection models."""

from __future__ import annotations

import io
import os
import re
import zipfile
from pathlib import Path
from typing import Any

DocumentSource = str | Path | bytes
DocumentStream = io.BufferedReader | io.BytesIO

DOCUMENT_MARKDOWN_CHAR_CAP = 4_000_000
_DOCUMENT_SOURCE_MAX_BYTES = 200 * 1024 * 1024
_OOXML_ENTRY_MAX_BYTES = 256 * 1024 * 1024
_OOXML_TOTAL_MAX_BYTES = 512 * 1024 * 1024
_OOXML_ENTRY_MAX_COUNT = 10_000
_OOXML_COMPRESSION_RATIO_MAX = 1_000
_XLSX_CELL_MAX_COUNT = 100_000

_OOXML_MARKERS = {
    ".docx": "word/document.xml",
    ".xlsx": "xl/workbook.xml",
    ".pptx": "ppt/presentation.xml",
}
_PARTIAL_NUMBERING = re.compile(r"^\.\d+$")


def _validate_source_size(stream: DocumentStream) -> None:
    try:
        size = os.fstat(stream.fileno()).st_size
    except (AttributeError, io.UnsupportedOperation):
        position = stream.tell()
        stream.seek(0, os.SEEK_END)
        size = stream.tell()
        stream.seek(position)
    if size > _DOCUMENT_SOURCE_MAX_BYTES:
        raise ValueError("document exceeds the 200 MB parsing limit")


def _validate_ooxml(stream: DocumentStream, extension: str) -> None:
    try:
        with zipfile.ZipFile(stream) as archive:
            entries = archive.infolist()
            names = {entry.filename for entry in entries}
            if "[Content_Types].xml" not in names or _OOXML_MARKERS[extension] not in names:
                raise ValueError(f"invalid {extension} package")
            if len(entries) > _OOXML_ENTRY_MAX_COUNT:
                raise ValueError(f"{extension} package has too many entries")
            total_size = 0
            for entry in entries:
                if entry.flag_bits & 1:
                    raise ValueError(f"encrypted {extension} packages are unsupported")
                if entry.file_size > _OOXML_ENTRY_MAX_BYTES:
                    raise ValueError(f"{extension} package entry is too large")
                if (
                    entry.file_size > 1024 * 1024
                    and entry.file_size > max(1, entry.compress_size) * _OOXML_COMPRESSION_RATIO_MAX
                ):
                    raise ValueError(f"{extension} package compression ratio is unsafe")
                total_size += entry.file_size
                if total_size > _OOXML_TOTAL_MAX_BYTES:
                    raise ValueError(f"{extension} package expands beyond the safe limit")
    except zipfile.BadZipFile as exc:
        raise ValueError(f"invalid {extension} package") from exc
    finally:
        stream.seek(0)


def _write_limited(output: io.StringIO, text: str, limit: int) -> bool:
    remaining = limit - output.tell()
    if remaining <= 0:
        return False
    output.write(text[:remaining])
    return len(text) <= remaining


def _write_block(output: io.StringIO, text: str, limit: int) -> bool:
    if not text:
        return True
    prefix = "\n\n" if output.tell() else ""
    return _write_limited(output, prefix + text, limit)


def _escape_cell(value: object) -> str:
    return str(value if value is not None else "").replace("|", "\\|").replace("\n", " ")


def _markdown_table(rows: list[list[Any]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    rendered = [[_escape_cell(cell) for cell in row] for row in normalized]
    return "\n".join(
        [
            "| " + " | ".join(rendered[0]) + " |",
            "| " + " | ".join("---" for _ in range(width)) + " |",
            *("| " + " | ".join(row) + " |" for row in rendered[1:]),
        ]
    )


def _docx_paragraph_markdown(paragraph: Any) -> str:
    from docx.text.hyperlink import Hyperlink

    parts: list[str] = []
    for item in paragraph.iter_inner_content():
        if isinstance(item, Hyperlink):
            target = item.url or (f"#{item.fragment}" if item.fragment else "")
            label = item.text.strip()
            parts.append(f"[{label}]({target.replace(')', r'\)')})" if target else label)
            continue
        text = item.text
        if text:
            if item.bold:
                text = f"**{text}**"
            if item.italic:
                text = f"*{text}*"
            parts.append(text)
        for properties in item._r.xpath(".//wp:docPr"):
            description = (
                properties.get("descr") or properties.get("title") or properties.get("name")
            )
            parts.append(f"[Image: {description or 'embedded image'}]")
    return "".join(parts).strip()


def _docx_markdown(stream: DocumentStream, limit: int) -> tuple[str, bool]:
    from docx import Document
    from docx.table import Table

    document = Document(stream)
    output = io.StringIO()
    for block in document.iter_inner_content():
        if isinstance(block, Table):
            rows = [[cell.text.strip() for cell in row.cells] for row in block.rows]
            if not _write_block(output, _markdown_table(rows), limit):
                return output.getvalue(), True
            continue
        text = _docx_paragraph_markdown(block)
        if not text:
            continue
        style = str(getattr(block.style, "name", "") or "")
        heading = re.fullmatch(r"Heading ([1-9])", style, flags=re.IGNORECASE)
        properties = block._p.pPr
        if heading:
            text = f"{'#' * int(heading.group(1))} {text}"
        elif style.startswith("List ") or (properties is not None and properties.numPr is not None):
            text = f"- {text}"
        if not _write_block(output, text, limit):
            return output.getvalue(), True
    return output.getvalue(), False


def _xlsx_markdown(stream: DocumentStream, limit: int) -> tuple[str, bool]:
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    workbook = load_workbook(stream, data_only=True, read_only=True, keep_links=False)
    output = io.StringIO()
    try:
        for sheet in workbook.worksheets:
            if not _write_block(output, f"## {sheet.title}", limit):
                return output.getvalue(), True
            max_column = max(1, min(int(sheet.max_column or 1), _XLSX_CELL_MAX_COUNT))
            max_row = max(
                1,
                min(
                    int(sheet.max_row or 1),
                    _XLSX_CELL_MAX_COUNT // max_column,
                ),
            )
            first_row = True
            for values in sheet.iter_rows(
                min_row=1,
                max_row=max_row,
                max_col=max_column,
                values_only=True,
            ):
                rendered = [_escape_cell(value) for value in values]
                while rendered and not rendered[-1]:
                    rendered.pop()
                if not rendered:
                    continue
                row = "| " + " | ".join(rendered) + " |"
                if not _write_limited(output, "\n" + row, limit):
                    return output.getvalue(), True
                if first_row:
                    separator = "| " + " | ".join("---" for _ in rendered) + " |"
                    if not _write_limited(output, "\n" + separator, limit):
                        return output.getvalue(), True
                    first_row = False
            if int(sheet.max_row or 0) * int(sheet.max_column or 0) > _XLSX_CELL_MAX_COUNT:
                return output.getvalue(), True
        return output.getvalue(), False
    finally:
        workbook.close()


def _chart_markdown(chart: Any) -> str:
    try:
        heading = "### Chart"
        if chart.has_title:
            heading += f": {chart.chart_title.text_frame.text}"
        categories = [str(category.label) for category in chart.plots[0].categories]
        series = list(chart.series)
        rows: list[list[object]] = [["Category", *(item.name for item in series)]]
        for index, category in enumerate(categories):
            rows.append([category, *(item.values[index] for item in series)])
        return f"{heading}\n\n{_markdown_table(rows)}"
    except (AttributeError, IndexError, TypeError, ValueError):
        return ""


def _pptx_markdown(stream: DocumentStream, limit: int) -> tuple[str, bool]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(stream)
    output = io.StringIO()

    def append_shape(shape: Any, title_element: Any | None) -> bool:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            return all(
                append_shape(child, title_element)
                for child in sorted(shape.shapes, key=lambda item: (item.top, item.left))
            )
        if getattr(shape, "has_table", False):
            rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
            return _write_block(output, _markdown_table(rows), limit)
        if getattr(shape, "has_chart", False):
            return _write_block(output, _chart_markdown(shape.chart), limit)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            description = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
            return _write_block(output, f"[Image: {description or shape.name}]", limit)
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            text = shape.text.strip()
            if title_element is not None and shape._element is title_element:
                text = f"# {text}"
            return _write_block(output, text, limit)
        return True

    for index, slide in enumerate(presentation.slides, start=1):
        if not _write_block(output, f"<!-- Slide number: {index} -->", limit):
            return output.getvalue(), True
        title = slide.shapes.title
        title_element = title._element if title is not None else None
        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            if not append_shape(shape, title_element):
                return output.getvalue(), True
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes and not _write_block(output, f"### Notes:\n\n{notes}", limit):
                return output.getvalue(), True
    return output.getvalue(), False


def _merge_partial_numbering(text: str) -> str:
    lines = text.splitlines()
    merged: list[str] = []
    index = 0
    while index < len(lines):
        current = lines[index].strip()
        if _PARTIAL_NUMBERING.fullmatch(current):
            next_index = index + 1
            while next_index < len(lines) and not lines[next_index].strip():
                next_index += 1
            if next_index < len(lines):
                merged.append(f"{current} {lines[next_index].strip()}")
                index = next_index + 1
                continue
        merged.append(lines[index])
        index += 1
    return "\n".join(merged)


def _pdf_markdown(stream: DocumentStream, limit: int) -> tuple[str, bool]:
    import pdfplumber

    header = stream.read(1024)
    if b"%PDF-" not in header:
        raise ValueError("invalid .pdf file")
    stream.seek(0)
    output = io.StringIO()
    with pdfplumber.open(stream) as document:
        for page_number, page in enumerate(document.pages, start=1):
            if not _write_block(output, f"<!-- Page number: {page_number} -->", limit):
                return output.getvalue(), True
            tables = page.find_tables()
            table_boxes = tuple(table.bbox for table in tables)

            def outside_tables(
                item: dict[str, Any],
                boxes: tuple[tuple[float, float, float, float], ...] = table_boxes,
            ) -> bool:
                if item.get("object_type") != "char":
                    return True
                horizontal = (float(item["x0"]) + float(item["x1"])) / 2
                vertical = (float(item["top"]) + float(item["bottom"])) / 2
                return all(
                    not (x0 <= horizontal <= x1 and top <= vertical <= bottom)
                    for x0, top, x1, bottom in boxes
                )

            text = _merge_partial_numbering(
                page.filter(outside_tables).extract_text(layout=True) or ""
            ).strip()
            if text and not _write_block(output, text, limit):
                return output.getvalue(), True
            for table in tables:
                if not _write_block(output, _markdown_table(table.extract()), limit):
                    return output.getvalue(), True
    return output.getvalue(), False


def _convert(stream: DocumentStream, extension: str, limit: int) -> tuple[str, bool]:
    if extension in _OOXML_MARKERS:
        _validate_ooxml(stream, extension)
    if extension == ".docx":
        return _docx_markdown(stream, limit)
    if extension == ".xlsx":
        return _xlsx_markdown(stream, limit)
    if extension == ".pptx":
        return _pptx_markdown(stream, limit)
    if extension == ".pdf":
        return _pdf_markdown(stream, limit)
    raise ValueError(f"unsupported document extension: {extension}")


def document_to_markdown(
    source: DocumentSource,
    extension: str,
    *,
    char_limit: int = DOCUMENT_MARKDOWN_CHAR_CAP,
) -> tuple[str, bool]:
    """Convert one bounded snapshot by its known suffix; return text and truncation."""
    if char_limit <= 0:
        raise ValueError("character limit must be positive")
    extension = extension.lower()
    if isinstance(source, bytes):
        byte_stream = io.BytesIO(source)
        _validate_source_size(byte_stream)
        return _convert(byte_stream, extension, char_limit)
    with Path(source).open("rb") as file_stream:
        _validate_source_size(file_stream)
        return _convert(file_stream, extension, char_limit)
